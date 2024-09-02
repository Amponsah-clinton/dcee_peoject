import os
from django.shortcuts import render, redirect
from django.contrib.auth.decorators import login_required
from django.contrib.auth import login, authenticate
from django.contrib.auth.forms import AuthenticationForm
from django.http import JsonResponse
from .forms import UserRegisterForm, PDFUploadForm
from .models import PDFDocument
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from datetime import datetime

load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

def get_pdf_text(pdf_files):
    text = ""
    for pdf in pdf_files:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=50000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")




import requests
from django.contrib.auth.decorators import login_required
from django.http import JsonResponse
from .models import PDFDocument, QuestionAnswer

def get_conversational_chain():
    prompt_template = """
    Answer the question as a complete sentence, providing as much detail as possible. If the answer is not in the provided context, just say, "The answer is not available in the context."
    Context:
    {context}
    Question:
    {question}
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain



import requests
from bs4 import BeautifulSoup
from django.contrib.auth.decorators import login_required
from django.http import JsonResponse
from .models import PDFDocument, QuestionAnswer
import requests
from bs4 import BeautifulSoup

def web_scrape_search(query):
    search_url = f"https://www.google.com/search?q={query.replace(' ', '+')}"
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3"
    }
    
    response = requests.get(search_url, headers=headers)
    soup = BeautifulSoup(response.text, "html.parser")
    snippet = None
    try:
        snippet = soup.find("span", {"class": "aCOpRe"})
        if snippet:
            return snippet.text
        else:
            snippet = soup.find("div", {"class": "BNeawe s3v9rd AP7Wnd"})
            if snippet:
                return snippet.text
    except AttributeError:
        pass

    return "Sorry, I couldn't find an answer to your question."




@login_required
def question_answer(request):
    if request.method == 'POST':
        user_question = request.POST.get('question')
        print("Received question:", user_question)
        qa_list = QuestionAnswer.objects.filter(user=request.user).order_by('-created_at')[:10]
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(user_question)

        chain = get_conversational_chain()
        response = chain(
            {"input_documents": docs, "question": user_question},
            return_only_outputs=True
        )
        answer = response["output_text"]
        if answer == "The answer is not available in the context.":
            answer = web_scrape_search(user_question)
            if not answer:
                answer = "Sorry, I couldn't find an answer to your question."
        print("Response:", answer)
        
        qa_entry = QuestionAnswer.objects.create(
            user=request.user,
            question=user_question,
            answer=answer
        )
        return JsonResponse({"response": answer})

    qa_list = QuestionAnswer.objects.filter(user=request.user).order_by('-created_at')[:10]
    return render(request, 'dcee_app/question_answer.html', {
        'qa_list': qa_list,
        'username': request.user.username 
    })



@login_required
def all_questions(request):
    all_qa_list = QuestionAnswer.objects.filter(user=request.user).order_by('-created_at')

    return render(request, 'dcee_app/all_questions.html', {
        'all_qa_list': all_qa_list,
    })


from docx import Document
from pptx import Presentation
import os

def get_file_text(file):
    extension = file.name.split('.')[-1].lower()
    
    if extension == 'pdf':
        return extract_text_from_pdf(file)
    elif extension == 'docx':
        return extract_text_from_docx(file)
    elif extension == 'pptx':
        return extract_text_from_pptx(file)
    elif extension == 'txt':
        return extract_text_from_txt(file)
    elif extension == 'ppt':
        return extract_text_from_txt(file)
    else:
        raise ValueError('Unsupported file type')

def extract_text_from_pdf(file):
    text = ""
    pdf_reader = PdfReader(file)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(file):
    doc = Document(file)
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_text_from_txt(file):
    return file.read().decode('utf-8')



@login_required
def pdf_upload(request):
    if request.method == 'POST':
        form = PDFUploadForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            file = request.FILES['file']
            raw_text = get_file_text(file)
            text_chunks = get_text_chunks(raw_text)
            get_vector_store(text_chunks)
            return redirect('pdf_upload')
    else:
        form = PDFUploadForm()
    return render(request, 'dcee_app/pdf_upload.html', {'form': form})





def register(request):
    if request.method == 'POST':
        form = UserRegisterForm(request.POST)
        if form.is_valid():
            form.save()
            username = form.cleaned_data.get('username')
            password = form.cleaned_data.get('password1')
            user = authenticate(username=username, password=password)
            login(request, user)
            return redirect('question_answer')
    else:
        form = UserRegisterForm()
    return render(request, 'dcee_app/register.html', {'form': form})




def login_view(request):
    if request.method == 'POST':
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            username = form.cleaned_data.get('username')
            password = form.cleaned_data.get('password')
            user = authenticate(username=username, password=password)
            if user is not None:
                login(request, user)
                return redirect('question_answer')
    else:
        form = AuthenticationForm()
    return render(request, 'dcee_app/login.html', {'form': form})



from django.http import JsonResponse
from django.views.decorators.http import require_POST
from .models import QuestionAnswer, Vote
import json
from django.views.decorators.http import require_http_methods

@require_http_methods(["PATCH"])
@login_required
def vote(request):
    try:
        data = json.loads(request.body)
        qa_id = data.get('qa_id')
        action = data.get('action')
        question_answer = QuestionAnswer.objects.get(id=qa_id)
        vote, created = Vote.objects.get_or_create(
            user=request.user, 
            question_answer=question_answer,
            defaults={'vote_type': action}
        )

        if not created: 
            if vote.vote_type == action:
                return JsonResponse({'error': 'You have already voted this way'}, status=400)
            else:
                vote.vote_type = action
                vote.save()

               
                if action == 'upvote':
                    question_answer.upvotes += 1
                    question_answer.downvotes -= 1
                else:
                    question_answer.upvotes -= 1
                    question_answer.downvotes += 1

        else:
            if action == 'upvote':
                question_answer.upvotes += 1
            elif action == 'downvote':
                question_answer.downvotes += 1

        question_answer.save()

        return JsonResponse({
            'upvotes': question_answer.upvotes,
            'downvotes': question_answer.downvotes
        })

    except QuestionAnswer.DoesNotExist:
        return JsonResponse({'error': 'QuestionAnswer not found'}, status=404)
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


from django.contrib.auth.forms import PasswordResetForm
from django.contrib.auth import get_user_model
from django.core.mail import send_mail
from django.template.loader import render_to_string
from django.utils.http import urlsafe_base64_encode
from django.utils.encoding import force_bytes
from django.contrib.auth.tokens import default_token_generator
from django.urls import reverse
from django.conf import settings

User = get_user_model()

def forgot_password_view(request):
    if request.method == 'POST':
        form = PasswordResetForm(request.POST)
        if form.is_valid():
            email = form.cleaned_data['email']
            associated_users = User.objects.filter(email=email)
            if associated_users.exists():
                for user in associated_users:
                    subject = "Password Reset Requested"
                    email_template_name = "password_reset_email.txt"
                    c = {
                        "email": user.email,
                        'domain': request.META['HTTP_HOST'],
                        'site_name': 'Your Site',
                        "uid": urlsafe_base64_encode(force_bytes(user.pk)),
                        "user": user,
                        'token': default_token_generator.make_token(user),
                        'protocol': 'http',
                    }
                    email = render_to_string(email_template_name, c)
                    send_mail(subject, email, settings.DEFAULT_FROM_EMAIL, [user.email], fail_silently=False)
            return redirect('password_reset_done')
    else:
        form = PasswordResetForm()
    return render(request, 'forgot_password.html', {'form': form})

def about(request):
    return render(request, 'about.html')
from django.contrib.auth import logout




from django.contrib.auth import logout

def logout_view(request):
    logout(request)
    return redirect('login')



def index(request):
    return render(request, 'index.html')


